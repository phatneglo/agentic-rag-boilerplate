"""
Document converter worker.
Processes document conversion jobs from the Redis queue using Marker library.
"""
import asyncio
import os
import sys
from pathlib import Path
from typing import Any, Dict
import signal

# Force CPU-only processing for Marker
os.environ["CUDA_VISIBLE_DEVICES"] = ""
os.environ["TORCH_DEVICE"] = "cpu"

# Add the parent directory to the Python path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from bullmq import Worker
from marker.converters.pdf import PdfConverter
from marker.models import create_model_dict
from marker.output import text_from_rendered

from app.core.config import settings
from app.core.logging_config import configure_logging, get_logger, log_job_event
from app.utils.exceptions import DocumentConversionError, FileProcessingError


# Configure logging
configure_logging()
logger = get_logger(__name__)


class DocumentConverterWorker:
    """Worker for processing document conversion jobs using Marker."""
    
    def __init__(self):
        self.worker = None
        self.marker_converter = None
        self.is_running = False
        self.shutdown_event = asyncio.Event()
    
    async def setup(self):
        """Setup worker and load Marker models."""
        try:
            # Initialize Marker converter (this may take some time on first run)
            logger.info("Loading Marker models...")
            
            # Force CPU usage for models
            import torch
            if torch.cuda.is_available():
                logger.info("CUDA is available but forcing CPU usage for compatibility")
            else:
                logger.info("CUDA not available, using CPU")
            
            # Create model dict with CPU device
            model_dict = create_model_dict()
            
            self.marker_converter = PdfConverter(
                artifact_dict=model_dict,
            )
            logger.info("Marker models loaded successfully")
            
            # Create Redis connection string
            if settings.redis_password:
                redis_connection = f"redis://:{settings.redis_password}@{settings.redis_host}:{settings.redis_port}/{settings.redis_db}"
            else:
                redis_connection = f"redis://{settings.redis_host}:{settings.redis_port}/{settings.redis_db}"
            
            # Create worker - BullMQ Python API
            self.worker = Worker(
                settings.queue_names["document_converter"],
                self.process_job,
                {"connection": redis_connection}
            )
            
            self.is_running = True
            
            logger.info(
                "Document converter worker initialized and started",
                queue_name=settings.queue_names["document_converter"]
            )
            
        except Exception as e:
            logger.error("Failed to setup document converter worker", error=str(e))
            raise
    
    async def process_job(self, job, job_token) -> Dict[str, Any]:
        """
        Process a document conversion job using Marker library.
        
        Args:
            job: BullMQ job object
            job_token: Job token for this processing instance
            
        Returns:
            Dict[str, Any]: Job result
            
        Raises:
            DocumentConversionError: If conversion fails
        """
        job_id = job.id
        job_data = job.data
        
        try:
            logger.info(
                "Processing document conversion job",
                job_id=job_id,
                document_id=job_data.get("document_id"),
                source_path=job_data.get("source_path")
            )
            
            log_job_event(
                job_id=job_id,
                queue_name=settings.queue_names["document_converter"],
                event_type="started",
                document_id=job_data.get("document_id")
            )
            
            # Extract job data
            document_id = job_data["document_id"]
            source_path = job_data["source_path"]
            output_path = job_data["output_path"]
            conversion_options = job_data.get("conversion_options", {})
            
            # Update job progress
            await job.updateProgress(10)
            
            # Validate source file exists
            if not os.path.exists(source_path):
                raise DocumentConversionError(f"Source file not found: {source_path}")
            
            # Ensure output directory exists
            output_dir = os.path.dirname(output_path)
            os.makedirs(output_dir, exist_ok=True)
            
            await job.updateProgress(20)
            
            # Convert document based on file type
            source_ext = os.path.splitext(source_path.lower())[1]
            
            if source_ext == '.pdf':
                result = await self._convert_pdf_with_marker(
                    source_path, output_path, conversion_options
                )
            elif source_ext in ['.docx', '.doc']:
                result = await self._convert_docx_to_markdown(
                    source_path, output_path, conversion_options
                )
            elif source_ext in ['.txt', '.md']:
                result = await self._convert_text_to_markdown(
                    source_path, output_path, conversion_options
                )
            elif source_ext in ['.html', '.htm']:
                result = await self._convert_html_to_markdown(
                    source_path, output_path, conversion_options
                )
            elif source_ext in ['.pptx', '.ppt']:
                result = await self._convert_pptx_with_marker(
                    source_path, output_path, conversion_options
                )
            elif source_ext in ['.xlsx', '.xls']:
                result = await self._convert_xlsx_with_marker(
                    source_path, output_path, conversion_options
                )
            elif source_ext == '.epub':
                result = await self._convert_epub_with_marker(
                    source_path, output_path, conversion_options
                )
            else:
                raise DocumentConversionError(f"Unsupported file format: {source_ext}")
            
            await job.updateProgress(90)
            
            # Prepare result
            job_result = {
                "success": True,
                "document_id": document_id,
                "source_path": source_path,
                "output_path": output_path,
                "conversion_result": result,
                "processed_at": job_data.get("created_at"),
            }
            
            await job.updateProgress(100)
            
            logger.info(
                "Document conversion job completed successfully",
                job_id=job_id,
                document_id=document_id,
                output_path=output_path
            )
            
            log_job_event(
                job_id=job_id,
                queue_name=settings.queue_names["document_converter"],
                event_type="completed",
                document_id=document_id,
                result=job_result
            )
            
            return job_result
            
        except Exception as e:
            logger.error(
                "Document conversion job failed",
                job_id=job_id,
                document_id=job_data.get("document_id"),
                error=str(e)
            )
            
            log_job_event(
                job_id=job_id,
                queue_name=settings.queue_names["document_converter"],
                event_type="failed",
                document_id=job_data.get("document_id"),
                error=str(e)
            )
            
            raise DocumentConversionError(f"Document conversion failed: {e}")

    async def _convert_pdf_with_marker(
        self,
        source_path: str,
        output_path: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert PDF to Markdown using Marker."""
        logger.info("Converting PDF to Markdown with Marker", source_path=source_path)
        
        try:
            # Run Marker conversion in thread pool to avoid blocking
            def convert_pdf():
                # Use the new Marker API
                rendered = self.marker_converter(source_path)
                text, metadata, images = text_from_rendered(rendered)
                return text, metadata, images
            
            loop = asyncio.get_event_loop()
            full_text, out_meta, images = await loop.run_in_executor(None, convert_pdf)
            
            # Save markdown content to output file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            
            # Save images if any (optional)
            if images and options.get("save_images", False):
                images_dir = os.path.join(os.path.dirname(output_path), "images")
                os.makedirs(images_dir, exist_ok=True)
                for filename, image_data in images.items():
                    image_path = os.path.join(images_dir, filename)
                    with open(image_path, 'wb') as f:
                        f.write(image_data)
            
            return {
                "format": "pdf",
                "pages_processed": out_meta.get("page_stats", []),
                "images_extracted": len(images) if images else 0,
                "output_size": len(full_text),
                "metadata": out_meta,
                "success": True
            }
            
        except Exception as e:
            logger.error("Marker PDF conversion failed", error=str(e))
            raise DocumentConversionError(f"PDF conversion with Marker failed: {e}")
    
    async def _convert_pptx_with_marker(
        self,
        source_path: str,
        output_path: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert PPTX to Markdown using Marker (if supported) or fallback."""
        logger.info("Converting PPTX to Markdown", source_path=source_path)
        
        try:
            # Try Marker first if it supports PPTX
            def convert_pptx():
                rendered = self.marker_converter(source_path)
                text, metadata, images = text_from_rendered(rendered)
                return text, metadata, images
            
            loop = asyncio.get_event_loop()
            full_text, out_meta, images = await loop.run_in_executor(None, convert_pptx)
            
            # Save markdown content
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            
            return {
                "format": "pptx", 
                "slides_processed": out_meta.get("page_stats", []),
                "images_extracted": len(images) if images else 0,
                "output_size": len(full_text),
                "metadata": out_meta,
                "success": True
            }
            
        except Exception as e:
            logger.error("Marker PPTX conversion failed, using fallback", error=str(e))
            # Fallback to basic text extraction if Marker fails
            return await self._fallback_pptx_conversion(source_path, output_path)
    
    async def _convert_xlsx_with_marker(
        self,
        source_path: str,
        output_path: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert XLSX to Markdown using Marker or fallback."""
        logger.info("Converting XLSX to Markdown", source_path=source_path)
        
        try:
            # Try Marker first if it supports XLSX
            def convert_xlsx():
                rendered = self.marker_converter(source_path)
                text, metadata, images = text_from_rendered(rendered)
                return text, metadata, images
            
            loop = asyncio.get_event_loop()
            full_text, out_meta, images = await loop.run_in_executor(None, convert_xlsx)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            
            return {
                "format": "xlsx",
                "sheets_processed": out_meta.get("page_stats", []),
                "output_size": len(full_text),
                "metadata": out_meta,
                "success": True
            }
            
        except Exception as e:
            logger.info("Marker XLSX conversion not supported, using fallback", error=str(e))
            return await self._fallback_xlsx_conversion(source_path, output_path)
    
    async def _convert_epub_with_marker(
        self,
        source_path: str,
        output_path: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert EPUB to Markdown using Marker."""
        logger.info("Converting EPUB to Markdown with Marker", source_path=source_path)
        
        try:
            def convert_epub():
                rendered = self.marker_converter(source_path)
                text, metadata, images = text_from_rendered(rendered)
                return text, metadata, images
            
            loop = asyncio.get_event_loop()
            full_text, out_meta, images = await loop.run_in_executor(None, convert_epub)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            
            return {
                "format": "epub",
                "chapters_processed": out_meta.get("page_stats", []),
                "output_size": len(full_text),
                "metadata": out_meta,
                "success": True
            }
            
        except Exception as e:
            logger.error("Marker EPUB conversion failed", error=str(e))
            return await self._fallback_epub_conversion(source_path, output_path)

    async def _convert_docx_to_markdown(
        self,
        source_path: str,
        output_path: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert DOCX to Markdown using python-docx."""
        logger.info("Converting DOCX to Markdown", source_path=source_path)
        
        try:
            from docx import Document
            import re
            
            def convert_docx():
                doc = Document(source_path)
                markdown_content = []
                
                for paragraph in doc.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Simple markdown formatting based on style
                        if paragraph.style.name.startswith('Heading'):
                            level = int(paragraph.style.name.split()[-1]) if paragraph.style.name.split()[-1].isdigit() else 1
                            markdown_content.append(f"{'#' * level} {text}\n")
                        else:
                            markdown_content.append(f"{text}\n")
                
                return '\n'.join(markdown_content)
            
            loop = asyncio.get_event_loop()
            markdown_text = await loop.run_in_executor(None, convert_docx)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_text)
            
            return {
                "format": "docx",
                "pages_processed": 1,
                "output_size": len(markdown_text),
                "success": True
            }
            
        except Exception as e:
            logger.error("DOCX conversion failed", error=str(e))
            raise DocumentConversionError(f"DOCX conversion failed: {e}")

    async def _convert_text_to_markdown(
        self,
        source_path: str,
        output_path: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert text file to Markdown."""
        logger.info("Converting text to Markdown", source_path=source_path)
        
        try:
            def convert_text():
                with open(source_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                # If it's already markdown, keep as is
                if source_path.lower().endswith('.md'):
                    return content
                
                # For plain text, add basic markdown formatting
                lines = content.split('\n')
                markdown_lines = []
                
                for line in lines:
                    line = line.strip()
                    if line:
                        # Simple heuristics for headers
                        if line.isupper() and len(line) < 80:
                            markdown_lines.append(f"## {line}")
                        else:
                            markdown_lines.append(line)
                    else:
                        markdown_lines.append("")
                
                return '\n'.join(markdown_lines)
            
            loop = asyncio.get_event_loop()
            markdown_content = await loop.run_in_executor(None, convert_text)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                "format": "text",
                "pages_processed": 1,
                "output_size": len(markdown_content),
                "success": True
            }
            
        except Exception as e:
            logger.error("Text conversion failed", error=str(e))
            raise DocumentConversionError(f"Text conversion failed: {e}")

    async def _convert_html_to_markdown(
        self,
        source_path: str,
        output_path: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert HTML to Markdown using html2text."""
        logger.info("Converting HTML to Markdown", source_path=source_path)
        
        try:
            import html2text
            
            def convert_html():
                h = html2text.HTML2Text()
                h.ignore_links = False
                h.ignore_images = False
                
                with open(source_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                return h.handle(html_content)
            
            loop = asyncio.get_event_loop()
            markdown_content = await loop.run_in_executor(None, convert_html)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                "format": "html",
                "pages_processed": 1,
                "output_size": len(markdown_content),
                "success": True
            }
            
        except Exception as e:
            logger.error("HTML conversion failed", error=str(e))
            raise DocumentConversionError(f"HTML conversion failed: {e}")

    async def _fallback_pptx_conversion(self, source_path: str, output_path: str) -> Dict[str, Any]:
        """Fallback PPTX conversion using python-pptx."""
        logger.info("Using fallback PPTX conversion", source_path=source_path)
        
        try:
            def extract_text():
                from pptx import Presentation
                
                prs = Presentation(source_path)
                text_content = []
                
                for i, slide in enumerate(prs.slides):
                    text_content.append(f"## Slide {i + 1}")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_content.append(shape.text)
                    
                    text_content.append("")  # Empty line between slides
                
                return '\n'.join(text_content)
            
            loop = asyncio.get_event_loop()
            markdown_content = await loop.run_in_executor(None, extract_text)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                "format": "pptx",
                "slides_processed": markdown_content.count("## Slide"),
                "output_size": len(markdown_content),
                "conversion_method": "fallback",
                "success": True
            }
            
        except Exception as e:
            logger.error("Fallback PPTX conversion failed", error=str(e))
            raise DocumentConversionError(f"PPTX conversion failed: {e}")

    async def _fallback_xlsx_conversion(self, source_path: str, output_path: str) -> Dict[str, Any]:
        """Fallback XLSX conversion using pandas."""
        logger.info("Using fallback XLSX conversion", source_path=source_path)
        
        try:
            def extract_data():
                import pandas as pd
                
                # Read all sheets
                excel_file = pd.ExcelFile(source_path)
                markdown_content = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(source_path, sheet_name=sheet_name)
                    
                    markdown_content.append(f"## {sheet_name}")
                    markdown_content.append("")
                    
                    # Convert DataFrame to markdown table
                    markdown_table = df.to_markdown(index=False)
                    markdown_content.append(markdown_table)
                    markdown_content.append("")
                
                return '\n'.join(markdown_content)
            
            loop = asyncio.get_event_loop()
            markdown_content = await loop.run_in_executor(None, extract_data)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                "format": "xlsx",
                "sheets_processed": markdown_content.count("## "),
                "output_size": len(markdown_content),
                "conversion_method": "fallback",
                "success": True
            }
            
        except Exception as e:
            logger.error("Fallback XLSX conversion failed", error=str(e))
            raise DocumentConversionError(f"XLSX conversion failed: {e}")

    async def _fallback_epub_conversion(self, source_path: str, output_path: str) -> Dict[str, Any]:
        """Fallback EPUB conversion using ebooklib."""
        logger.info("Using fallback EPUB conversion", source_path=source_path)
        
        try:
            def extract_text():
                import ebooklib
                from ebooklib import epub
                from bs4 import BeautifulSoup
                
                book = epub.read_epub(source_path)
                chapters = []
                
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), 'html.parser')
                        text = soup.get_text()
                        if text.strip():
                            chapters.append(f"# Chapter\n\n{text}")
                
                return '\n\n---\n\n'.join(chapters)
            
            loop = asyncio.get_event_loop()
            markdown_content = await loop.run_in_executor(None, extract_text)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                "format": "epub",
                "chapters_processed": markdown_content.count("# Chapter"),
                "output_size": len(markdown_content),
                "conversion_method": "fallback",
                "success": True
            }
            
        except Exception as e:
            logger.error("Fallback EPUB conversion failed", error=str(e))
            raise DocumentConversionError(f"EPUB conversion failed: {e}")

    async def stop(self):
        """Stop the worker gracefully."""
        logger.info("Stopping document converter worker...")
        self.is_running = False
        self.shutdown_event.set()
        
        if self.worker:
            await self.worker.close()
            logger.info("Worker stopped gracefully")

    def signal_handler(self, signum, frame):
        """Handle shutdown signals."""
        logger.info("Received signal to stop document converter worker")
        asyncio.create_task(self.stop())


async def main():
    """Main function to run the document converter worker."""
    worker = DocumentConverterWorker()
    
    try:
        await worker.setup()
        
        # Set up signal handlers for graceful shutdown
        def signal_handler(signum, frame):
            logger.info("Signal received, shutting down.")
            worker.shutdown_event.set()
        
        # Assign signal handlers to SIGTERM and SIGINT
        signal.signal(signal.SIGTERM, signal_handler)
        signal.signal(signal.SIGINT, signal_handler)
        
        logger.info("Document converter worker is running. Press Ctrl+C to stop.")
        
        # Wait until the shutdown event is set
        await worker.shutdown_event.wait()
        
    except KeyboardInterrupt:
        logger.info("Document converter worker stopped by user")
    except Exception as e:
        logger.error(f"Document converter worker error: {e}")
    finally:
        # Close the worker
        logger.info("Cleaning up worker...")
        if worker.worker:
            await worker.worker.close()
        logger.info("Worker shut down successfully.")


if __name__ == "__main__":
    asyncio.run(main()) 